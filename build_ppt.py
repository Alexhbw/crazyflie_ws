#!/usr/bin/env python3
"""生成七八九十章 PPT 幻灯片"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUTPUT = "/home/alex/crazyflie_ws/安全监视与实验框架_PPT.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x00, 0x33, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x60, 0x60, 0x60)
RED = RGBColor(0xCC, 0x33, 0x33)
GREEN = RGBColor(0x33, 0x99, 0x33)
LIGHT_BG = RGBColor(0xF2, 0xF5, 0xF9)


def slide(title_text, subtitle=None):
    """创建带标题的空白页"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 顶部色条
    bar = sl.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.0))  # MSO_SHAPE.RECTANGLE=1
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    # 标题
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tb2 = sl.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(0.5))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
    return sl


def textbox(sl, left, top, width, height, text_lines, font_size=14, bold_first=False, color=None):
    """多行文本框"""
    tb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(6)
        if bold_first and i == 0:
            p.font.bold = True
    return tb


def table(sl, left, top, width, height, headers, rows, col_widths=None):
    """添加表格"""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    tbl = sl.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height)).table

    for ci, h in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = DARK
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.text = str(val)
            if ri % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT_BG
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER

    if col_widths:
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = Inches(w)
    return tbl


def img_placeholder(sl, left, top, width, height, caption):
    """图片占位框"""
    shp = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))  # rect
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    shp.line.color.rgb = GRAY
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[ {caption} ]"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    # caption below
    c = sl.shapes.add_textbox(Inches(left), Inches(top + height + 0.05), Inches(width), Inches(0.3))
    c.text_frame.paragraphs[0].text = caption
    c.text_frame.paragraphs[0].font.size = Pt(10)
    c.text_frame.paragraphs[0].font.color.rgb = GRAY
    c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# ========================================================================
# 封面
# ========================================================================
sl = slide("安全监视器与求解器实验框架", "面向无人机安全自主飞行的 ODE 稳定域分析与安全边界形式化验证")
textbox(sl, 1, 3.2, 11, 2.0, [
    "Crazyflie 2.1 集群仿真系统增强",
    "",
    "运行时安全监视器 (SafetyMonitor)",
    "区间可达集近似分析 (Interval Reachability)",
    "求解器稳定域对比实验框架",
    "自动化图表生成与报告系统",
], font_size=16, bold_first=True)
textbox(sl, 1, 6.5, 5, 0.5, ["2026 年 6 月"], font_size=12, color=GRAY)

# ========================================================================
# 总览
# ========================================================================
sl = slide("新增功能总览", "四大模块，覆盖安全监视 → 形式化验证 → 实验对比 → 自动化报告")

textbox(sl, 1, 1.4, 11, 0.6, ["📊  四大新增功能模块"], font_size=20, bold_first=True)

table(sl, 0.8, 2.2, 11.5, 4.0,
      ["模块", "核心功能", "关键指标/方法", "输出"],
      [
          ["安全监视器\n(SafetyMonitor)", "每步形式化安全条件检查\n碰撞/速度/姿态/能量漂移", "4 维 □不变式\n预热机制 (2s)\n安全分数体系", "safety_log.csv\nSafety Summary"],
          ["区间可达集\n(Interval Reachability)", "一步正向区间传播\n最坏情况碰撞预警", "采样法 (25次/步)\nJacobian法 (13次/步)\nAABB-圆柱碰撞检测", "实时 WARN 日志\nd_min 输出"],
          ["实验框架\n(Experiment Framework)", "全因子参数扫描\n临界步长二分搜索", "5求解器 × 5步长 × 3飞机数 × 2障碍物\n150次实验", "comparison_results.csv\nsafety_boundary_table.csv"],
          ["图表生成\n(Plot Results)", "自适应单/多步长\n柱状图/折线图/雷达图", "6 张报告级 PNG\n150 DPI 高清", "plots/*.png"],
      ],
)

# ========================================================================
# 第7章：SafetyMonitor
# ========================================================================
sl = slide("七. 运行时安全监视器 (SafetyMonitor)", "形式化安全条件的在线运行时监控")

textbox(sl, 1, 1.4, 5.5, 4.5, [
    "设计动机",
    "• 传统仿真依赖事后分析，无法实时检测违规",
    "• 每步自动检查预定义的安全条件",
    "• 将条件表达为时态逻辑不变式 □(condition)",
    "",
    "四维安全条件",
    "1. 碰撞安全: □(min d_surf > 0)  — d_surf > 0",
    "2. 速度安全: □(|v_i| < v_max)  — v_max = 2.5 m/s",
    "3. 姿态安全: □(|θ_i| < θ_max)  — θ_max = 0.8 rad",
    "4. 能量漂移: □(|ΔE/Δt| < ε)  — ε = 0.01 J/步 (诊断)",
], font_size=13)

textbox(sl, 7, 1.4, 5.5, 4.5, [
    "能量计算",
    "E = ½·m·|v|² + ½·ωᵀ·I·ω + m·g·z",
    "",
    "安全分数",
    "compliance = N_compliant / N_total",
    "overall = min(碰撞, 速度, 姿态)",
    "",
    "预热机制",
    "safety.warmup_time = 2.0 s",
    "启动初期跳过碰撞/姿态检查",
    "（无人机初始位置可能靠近障碍物）",
], font_size=13)

table(sl, 1, 6, 11.5, 1.2,
      ["参数", "safety.v_max", "safety.theta_max", "safety.warmup_time", "safety.log_path", "safety.verbose"],
      [
          ["默认值", "2.5", "0.8", "2.0", "safety_log.csv", "false"],
      ],
)

# ========================================================================
# 第8章：Interval Reachability
# ========================================================================
sl = slide("八. 区间可达集近似分析 (Interval Reachability)", "一步正向区间传播 — 最坏情况碰撞预警")

textbox(sl, 1, 1.4, 6, 5.0, [
    "问题定义",
    "• 对初始状态 x 施加 ±ε 有界扰动",
    "• 计算一步积分后的状态区间包围盒",
    "• 保守检测包围盒是否与障碍物相交",
    "• 在 ±ε 扰动下无交集 → 绝对不可能碰撞",
    "",
    "采样法 (默认, 保守性强)",
    "∀j ∈ [0,11]: x⁺ⱼ = solver(x + ε·eⱼ, u, h)",
    "            x⁻ⱼ = solver(x − ε·eⱼ, u, h)",
    "共 25 次求解器调用 (1 中心 + 12维×2方向)",
], font_size=13)

textbox(sl, 7.5, 1.4, 5.5, 5.0, [
    "Jacobian 法 (可选, 更快速)",
    "J[:,j] = [solver(+ε) − solver(−ε)] / (2ε)",
    "lower = x_center − |J|·ε·1₁₂",
    "upper = x_center + |J|·ε·1₁₂",
    "共 13 次求解器调用 (线性化近似)",
    "",
    "AABB-圆柱碰撞检测",
    "• 计算位置区间包围盒到圆柱轴心的最近点",
    "• d_xy < r_obs 且 |dz| < h/2 → 重叠告警",
    "",
    "参数: reachability.enabled / epsilon / method",
], font_size=13)

img_placeholder(sl, 1, 6.2, 6.0, 1.2, "区间可达集示意图 (可达包围盒 vs 障碍物圆柱)")

# ========================================================================
# 第9章：实验框架 — 实验矩阵
# ========================================================================
sl = slide("九. 求解器对比实验框架", "全因子参数扫描 — 150 次自动实验")

textbox(sl, 1, 1.4, 11, 0.6, ["实验设计: 全因子参数扫描矩阵"], font_size=18, bold_first=True)

table(sl, 0.8, 2.2, 11.5, 2.5,
      ["扫描维度", "取值集合", "数量", "说明"],
      [
          ["求解器", "Euler, Heun, RK4, RK45, Implicit Euler", "5", "一阶→四阶→自适应→A-稳定"],
          ["积分步长 (s)", "0.001, 0.005, 0.01, 0.02, 0.05", "5", "覆盖精密→粗放范围"],
          ["无人机数量", "1, 4, 8", "3", "单机→小队→集群"],
          ["障碍物", "true, false", "2", "有无环境障碍物"],
          ["仿真时长", "30.0 s (固定)", "1", ""],
      ],
)

textbox(sl, 1, 5.0, 11, 1.5, [
    "总计 5 × 5 × 3 × 2 = 150 次独立实验 | 每次 30s 仿真 | 全自动运行/采集/退出",
    "",
    "实验模式 (Headless): 跳过所有 ROS 发布 (TF/odom/marker/clock)，纯数值积分 + 指标采集，最大化运行效率",
], font_size=14)

table(sl, 0.8, 7, 11.5, 0.3,
      ["参数", "experiment.enabled", "experiment.duration", "experiment.headless", "experiment.metrics_path"],
      [
          ["默认值", "false", "30.0", "false", "experiment_metrics.csv"],
      ],
)

# ========================================================================
# 第9章：六维评估指标
# ========================================================================
sl = slide("九. 六维评估指标体系", "稳定性 · 安全性 · 精确性 · 实时性")

table(sl, 0.8, 1.6, 11.5, 3.5,
      ["维度", "指标", "符号/单位", "含义", "期望方向"],
      [
          ["稳定性", "能量漂移率", "dE/dt (J/s)", "总机械能变化速率，度量数值稳定性", "↓ 越小越稳定"],
          ["安全性", "最小障碍物距离", "d_min (m)", "全程最小 d_surf，负值 = 穿透", "↑ 越大越安全"],
          ["可控性", "最大线速度", "V_max (m/s)", "速度极值，>2.5 表失控", "↓ 越小越可控"],
          ["可控性", "最大姿态角", "θ_max (rad)", "倾斜极值，>0.8 表翻车风险", "↓ 越小越稳定"],
          ["精确性", "RMS 位置误差", "e_rms (m)", "相对 RK45+dt=0.001 参考轨迹", "↓ 越小越精确"],
          ["实时性", "平均计算耗时", "T_avg (ms)", "单步积分时间开销", "↓ 越小越实时"],
      ],
)

textbox(sl, 1, 5.5, 11, 1.2, [
    "核心发现:  安全步长 < 理论稳定步长",
    "数值稳定只是安全飞行的必要条件，而非充分条件。安全边界更窄，因为安全还需考虑物理约束（障碍物距离、速度上限、姿态限制）。",
], font_size=16, bold_first=True, color=RED)

# ========================================================================
# 第9章：临界步长 + 图表
# ========================================================================
sl = slide("九. 临界步长二分搜索 + 图表生成", "从实验数据到报告级可视化")

textbox(sl, 1, 1.4, 5.5, 2.5, [
    "临界步长二分搜索",
    "• 搜索范围: [0.001, 1.0] s, 精度 0.001 s",
    "• 安全条件: 无 COLLISION / VELOCITY / ATTITUDE",
    "    / ENERGY_DRIFT (>50 J) 违规",
    "• 输出: safe_step_size + unsafe_step_size + margin",
    "",
    "预期结果:",
    "Euler ≈ 0.02s  |  Heun ≈ 0.03s  |  RK4 ≈ 0.05s",
    "RK45 ≈ 0.05s  |  Implicit ≈ 1.0s (A-稳定)",
], font_size=13)

textbox(sl, 7, 1.4, 5.5, 2.5, [
    "图表流水线 (plot_results.py)",
    "• 自适应单/多步长: 柱状图 vs 折线图",
    "• 支持 quick_compare 和 full_sweep 两种格式",
    "• 6 张 150 DPI 报告级 PNG",
    "",
    "01  能量漂移 vs 步长",
    "02  障碍物距离 vs 步长",
    "03  跟踪误差 vs 步长",
    "04  计算耗时 vs 步长",
    "05  各求解器安全临界步长",
    "06  五维综合雷达图",
], font_size=13)

img_placeholder(sl, 1, 4.3, 5.8, 2.6, "01~03 实验对比图 (稳定性/安全性/精度)")
img_placeholder(sl, 7.2, 4.3, 5.8, 2.6, "04~06 实验对比图 (性能/安全边界/雷达)")

# ========================================================================
# 第10章：运行模式
# ========================================================================
sl = slide("十. 运行模式与使用指南", "7 种模式覆盖日常开发 → 批量实验 → 报告产出")

table(sl, 0.5, 1.6, 12.3, 3.5,
      ["模式", "命令", "GUI", "耗时", "用途"],
      [
          ["正常仿真", "./run_experiment.sh rk4 0.01 5", "RViz ✓", "实时", "日常开发与演示"],
          ["安全监视", "./run_safety_monitor.sh rk4 0.01 4", "RViz ✓", "实时", "违规实时告警"],
          ["单次采集", "./run_collect_data.sh rk4 0.01 4 60 true", "✗", "60s", "深度定量分析"],
          ["快速对比", "./run_quick_compare.sh", "✗", "~2min", "5求解器 × 15s"],
          ["多步长对比", "./run_multi_dt_compare.sh", "✗", "~5min", "5求解器 × 3步长"],
          ["全因子扫描", "python3 run_experiments.py", "✗", "~2h", "150次实验"],
          ["临界步长", "python3 safety_boundary_analysis.py", "✗", "~30min", "二分搜索安全边界"],
      ],
)

textbox(sl, 0.5, 5.3, 5.5, 1.5, [
    "典型工作流",
    "1. ./run_quick_compare.sh      ← 快速摸底",
    "2. ./run_multi_dt_compare.sh   ← 看步长影响",
    "3. python3 safety_boundary_analysis.py ← 找安全边界",
    "4. python3 plot_results.py [csv]        ← 出图",
], font_size=13, bold_first=True)

textbox(sl, 6.5, 5.3, 6.5, 1.5, [
    "关键参数速查",
    "solver: euler | heun | rk4 | rk45 | implicit",
    "step_size: 0.01 (正常) ~ 1.0 (极限测试)",
    "enable_obstacles: true | false",
    "experiment.enabled: true (自动采集模式)",
    "experiment.headless: true (无GUI, 加速)",
], font_size=13, bold_first=True)

# ========================================================================
# 总结
# ========================================================================
sl = slide("总结", "四大模块 → 完整的安全验证与实验链路")

textbox(sl, 1, 1.8, 5.5, 4.5, [
    "🔒 安全监视器",
    "• 四维 □不变式 每步检查",
    "• 安全分数自动聚合",
    "• CSV 违规日志 + 关机摘要",
    "",
    "📐 区间可达集",
    "• 采样法/Jacobian 法",
    "• 最坏情况碰撞预警",
    "• 轻量级，无外部依赖",
], font_size=14, bold_first=False)

textbox(sl, 7, 1.8, 5.5, 4.5, [
    "🔬 实验框架",
    "• 150 次全因子扫描",
    "• 临界步长二分搜索",
    "• 六维评估指标体系",
    "",
    "📊 图表生成",
    "• 自适应单/多步长图表",
    "• 6 张报告级 PNG",
    "• 直接插入论文/PPT",
], font_size=14, bold_first=False)

textbox(sl, 1, 6.2, 12, 1.0, [
    "技术贡献: 将数值分析 (ODE 稳定域) 与形式化方法 (安全条件不变式、区间可达集) 通过工程实现 (ROS 2 + C++ + Python) 有机融合，",
    "为「面向无人机安全自主飞行的 ODE 初值问题数值解法稳定域分析与安全边界形式化验证」提供了完整的实验验证平台。",
], font_size=13, color=RED)

# ========================================================================
# 保存
# ========================================================================
prs.save(OUTPUT)
size_kb = os.path.getsize(OUTPUT) / 1024
print(f"PPT 已生成: {OUTPUT}")
print(f"大小: {size_kb:.1f} KB")
print(f"幻灯片数: {len(prs.slides)}")
